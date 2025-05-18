from pptx import Presentation

def generate_ppt(slide_data):
    prs = Presentation()
    for title, content in slide_data:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = title
        slide.placeholders[1].text = content
    output_path = "AI2Slides_Output.pptx"
    prs.save(output_path)
    return output_path